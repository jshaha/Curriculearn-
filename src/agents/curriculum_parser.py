"""
Agent 1: Curriculum Parser

Extracts and structures lesson content from various document formats.

Supported formats:
- PDF (.pdf) - lesson plans, worksheets, reading passages
- PowerPoint (.pptx) - presentations
- Word (.docx) - documents
- Plain text (.txt) - transcripts, passages

Input: File path or raw text
Output: Structured lesson content for Agent 2
"""

import re
from pathlib import Path
from typing import List, Dict, Any, Optional
import pypdf
from pptx import Presentation
from docx import Document


class CurriculumParser:
    """
    Parses educational content from various file formats.

    Extracts text and intelligently segments it into lesson components
    for brain simulation and analysis.
    """

    def __init__(self):
        """Initialize the curriculum parser."""
        self.supported_formats = ['.pdf', '.pptx', '.docx', '.txt']

    def parse(self, file_path: str) -> Dict[str, Any]:
        """
        Parse a lesson file and extract structured content.

        Args:
            file_path: Path to the lesson file

        Returns:
            Dictionary containing:
                - 'sections': List of lesson segments
                - 'metadata': File information
                - 'source_type': Type of source document
        """
        path = Path(file_path)

        if not path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")

        extension = path.suffix.lower()

        if extension not in self.supported_formats:
            raise ValueError(
                f"Unsupported file format: {extension}. "
                f"Supported formats: {', '.join(self.supported_formats)}"
            )

        print(f"Parsing {extension} file: {path.name}")

        # Parse based on file type
        if extension == '.pdf':
            sections = self._parse_pdf(file_path)
            source_type = 'pdf'
        elif extension == '.pptx':
            sections = self._parse_pptx(file_path)
            source_type = 'powerpoint'
        elif extension == '.docx':
            sections = self._parse_docx(file_path)
            source_type = 'word'
        elif extension == '.txt':
            sections = self._parse_txt(file_path)
            source_type = 'text'

        # Clean and validate sections
        sections = self._clean_sections(sections)

        if not sections:
            raise ValueError(f"No content extracted from {file_path}")

        print(f"Extracted {len(sections)} lesson segments")

        return {
            'sections': sections,
            'metadata': {
                'filename': path.name,
                'format': extension,
                'num_sections': len(sections)
            },
            'source_type': source_type
        }

    def parse_text(self, text: str, segment_by: str = 'paragraph') -> Dict[str, Any]:
        """
        Parse raw text directly (useful for transcripts, pasted content).

        Args:
            text: Raw text content
            segment_by: How to segment ('paragraph', 'sentence', 'line')

        Returns:
            Structured lesson content
        """
        print(f"Parsing raw text ({len(text)} characters)")

        if segment_by == 'paragraph':
            sections = self._segment_by_paragraph(text)
        elif segment_by == 'sentence':
            sections = self._segment_by_sentence(text)
        elif segment_by == 'line':
            sections = text.split('\n')
        else:
            raise ValueError(f"Unknown segmentation method: {segment_by}")

        sections = self._clean_sections(sections)

        print(f"Extracted {len(sections)} segments")

        return {
            'sections': sections,
            'metadata': {
                'source': 'raw_text',
                'num_sections': len(sections),
                'segmentation': segment_by
            },
            'source_type': 'text'
        }

    def _parse_pdf(self, file_path: str) -> List[str]:
        """Extract text from PDF file."""
        sections = []

        try:
            reader = pypdf.PdfReader(file_path)
            num_pages = len(reader.pages)

            print(f"  Found {num_pages} pages")

            for page_num, page in enumerate(reader.pages):
                text = page.extract_text()

                if text.strip():
                    # Try to split by logical sections (paragraphs)
                    page_sections = self._segment_by_paragraph(text)

                    # Add page number context
                    for section in page_sections:
                        sections.append({
                            'content': section,
                            'page': page_num + 1
                        })

        except Exception as e:
            raise ValueError(f"Error parsing PDF: {str(e)}")

        return sections

    def _parse_pptx(self, file_path: str) -> List[str]:
        """Extract text from PowerPoint presentation."""
        sections = []

        try:
            prs = Presentation(file_path)
            num_slides = len(prs.slides)

            print(f"  Found {num_slides} slides")

            for slide_num, slide in enumerate(prs.slides):
                slide_text = []

                # Extract text from all shapes (title, content, etc.)
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())

                # Combine slide content
                if slide_text:
                    combined_text = " ".join(slide_text)
                    sections.append({
                        'content': combined_text,
                        'slide': slide_num + 1
                    })

        except Exception as e:
            raise ValueError(f"Error parsing PowerPoint: {str(e)}")

        return sections

    def _parse_docx(self, file_path: str) -> List[str]:
        """Extract text from Word document."""
        sections = []

        try:
            doc = Document(file_path)
            num_paragraphs = len(doc.paragraphs)

            print(f"  Found {num_paragraphs} paragraphs")

            for para_num, para in enumerate(doc.paragraphs):
                text = para.text.strip()

                if text:
                    # Check if it's a heading (larger sections)
                    if para.style.name.startswith('Heading'):
                        sections.append({
                            'content': text,
                            'type': 'heading',
                            'paragraph': para_num + 1
                        })
                    else:
                        sections.append({
                            'content': text,
                            'type': 'paragraph',
                            'paragraph': para_num + 1
                        })

        except Exception as e:
            raise ValueError(f"Error parsing Word document: {str(e)}")

        return sections

    def _parse_txt(self, file_path: str) -> List[str]:
        """Extract text from plain text file."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()

            # Segment by paragraphs
            sections = self._segment_by_paragraph(text)

            # Add line numbers
            result = []
            for i, section in enumerate(sections):
                result.append({
                    'content': section,
                    'segment': i + 1
                })

            return result

        except Exception as e:
            raise ValueError(f"Error parsing text file: {str(e)}")

    def _segment_by_paragraph(self, text: str) -> List[str]:
        """
        Segment text by paragraphs.

        Paragraphs are separated by blank lines or multiple newlines.
        """
        # Split by double newlines or more
        paragraphs = re.split(r'\n\s*\n', text)

        # Clean up
        paragraphs = [p.strip() for p in paragraphs if p.strip()]

        return paragraphs

    def _segment_by_sentence(self, text: str) -> List[str]:
        """
        Segment text by sentences.

        Uses simple sentence boundary detection.
        """
        # Simple sentence splitting (can be improved with nltk if needed)
        sentences = re.split(r'[.!?]+\s+', text)

        # Clean up
        sentences = [s.strip() for s in sentences if s.strip()]

        return sentences

    def _clean_sections(self, sections: List[Any]) -> List[str]:
        """
        Clean and normalize section text.

        Args:
            sections: List of sections (can be strings or dicts)

        Returns:
            List of clean text strings
        """
        cleaned = []

        for section in sections:
            # Extract text from dict or use directly
            if isinstance(section, dict):
                text = section.get('content', '')
            else:
                text = str(section)

            # Clean whitespace
            text = ' '.join(text.split())

            # Skip empty or very short sections
            if len(text) < 10:
                continue

            # Skip common noise patterns
            if self._is_noise(text):
                continue

            cleaned.append(text)

        return cleaned

    def _is_noise(self, text: str) -> bool:
        """
        Check if text segment is likely noise/metadata.

        Returns True if text should be filtered out.
        """
        # Filter out page numbers, headers, footers, etc.
        noise_patterns = [
            r'^Page \d+$',
            r'^\d+$',
            r'^[A-Z\s]{2,}$',  # ALL CAPS headers
            r'^Slide \d+$',
        ]

        for pattern in noise_patterns:
            if re.match(pattern, text, re.IGNORECASE):
                return True

        return False

    def get_preview(self, lesson_content: Dict[str, Any], max_sections: int = 5) -> str:
        """
        Generate a preview of parsed content.

        Args:
            lesson_content: Output from parse()
            max_sections: Maximum number of sections to show

        Returns:
            Formatted preview string
        """
        sections = lesson_content['sections']
        metadata = lesson_content['metadata']

        preview = f"\n{'='*60}\n"
        preview += f"CURRICULUM PARSER - PREVIEW\n"
        preview += f"{'='*60}\n"
        preview += f"File: {metadata.get('filename', 'N/A')}\n"
        preview += f"Format: {metadata.get('format', 'N/A')}\n"
        preview += f"Total Sections: {len(sections)}\n"
        preview += f"{'='*60}\n\n"

        for i, section in enumerate(sections[:max_sections]):
            preview += f"Section {i+1}:\n"
            preview += f"  {section[:150]}{'...' if len(section) > 150 else ''}\n\n"

        if len(sections) > max_sections:
            preview += f"... and {len(sections) - max_sections} more sections\n"

        preview += f"{'='*60}\n"

        return preview


# Convenience function for quick parsing
def parse_lesson(file_path: str) -> Dict[str, Any]:
    """
    Quick function to parse a lesson file.

    Args:
        file_path: Path to lesson file

    Returns:
        Structured lesson content
    """
    parser = CurriculumParser()
    return parser.parse(file_path)
